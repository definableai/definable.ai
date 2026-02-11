"""PPTX parser — extracts text from PowerPoint presentations using python-pptx.

Requires optional dependency: ``python-pptx>=1.0.0``
"""

import io
from typing import List, Set

from definable.readers.models import ContentBlock, ReaderConfig
from definable.readers.parsers.base_parser import BaseParser


def _import_pptx():
  try:
    from pptx import Presentation

    return Presentation
  except ImportError as e:
    raise ImportError(
      "PptxParser requires the 'python-pptx' package. Install it with: pip install 'python-pptx>=1.0.0' or: pip install 'definable-ai[readers]'"
    ) from e


class PptxParser(BaseParser):
  """Parses PPTX files using python-pptx (local, no API calls).

  Extracts text from all slides, one ContentBlock per slide.
  """

  def __init__(self) -> None:
    _import_pptx()

  def supported_mime_types(self) -> List[str]:
    return ["application/vnd.openxmlformats-officedocument.presentationml.presentation"]

  def supported_extensions(self) -> Set[str]:
    return {".pptx"}

  def parse(self, data: bytes, *, mime_type: str | None = None, config: ReaderConfig | None = None) -> List[ContentBlock]:
    Presentation = _import_pptx()
    prs = Presentation(io.BytesIO(data))
    blocks: List[ContentBlock] = []

    for i, slide in enumerate(prs.slides):
      parts: List[str] = []
      for shape in slide.shapes:
        if shape.has_text_frame:
          for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
              parts.append(text)
        if shape.has_table:
          table = shape.table
          rows = []
          for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("\t".join(cells))
          if rows:
            parts.append("\n".join(rows))

      if parts:
        blocks.append(
          ContentBlock(
            content_type="text",
            content="\n\n".join(parts),
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            page_number=i + 1,
            metadata={"slide_number": i + 1},
          )
        )

    return blocks
